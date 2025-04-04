import numpy as np
import matplotlib.pyplot as plt
from scipy.stats import norm
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.lang import MSO_LANGUAGE_ID
from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement

# функция для форматирования заголовка раздела/презентации
def presentation_title(tf_title):
    text_frame = tf_title.text_frame
    for run in text_frame.paragraphs[0].runs:
        run.font.size = Pt(40)
        run.font.bold = True
        run.font.name = "Cambria Math"
        run.font.language_id = MSO_LANGUAGE_ID.RUSSIAN

# функция для форматирования заголовка слайда
def title_format(tf_title):
    tf_title = slide.shapes.title.text_frame
    for run in tf_title.paragraphs[0].runs:
        run.font.size = Pt(28)
        run.font.name = "Cambria Math"
        run.font.bold = True
        run.font.language_id = MSO_LANGUAGE_ID.RUSSIAN


# функция для добавления текста слайда (placeholders)
def format_placeholders(theory, paragraph):
    theory.text = paragraph
    tf = theory.text_frame
    for run in tf.paragraphs[0].runs:
        run.font.size = Pt(24)
        run.font.name = "Arial"
        run.font.language_id = MSO_LANGUAGE_ID.RUSSIAN


# функция для форматирования абзаца
def add_paragraph_format(slide, paragraph):
    tf = slide.text_frame
    p = tf.add_paragraph()
    p.text = paragraph
    p.font.size = Pt(24)
    p.font.name = 'Arial'
    p.font.language_id = MSO_LANGUAGE_ID.RUSSIAN


# функция для форматирования абзаца с отступом
def add_indent_paragraph(theory, thesis):
    tf = theory.text_frame
    p = tf.add_paragraph()
    p.text = thesis
    p.font.size = Pt(24)
    p.font.name = 'Arial'
    p.level = 2
    p.font.language_id = MSO_LANGUAGE_ID.RUSSIAN


# функция для добавления надписи
def add_placeholder(paragraph, left, top, width, height):
    txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.text = paragraph
    for run in tf.paragraphs[0].runs:
        run.font.size = Pt(22)
        run.font.name = 'Arial'
        run.font.language_id = MSO_LANGUAGE_ID.RUSSIAN

# функция для выравнивания текста по центру
def paragraph_centered(paragraph, left, top, width, height):
    txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.line_spacing = 1.5
    tf.text = paragraph
    for run in tf.paragraphs[0].runs:
        run.font.size = Pt(22)
        run.font.name = 'Arial'
        run.font.language_id = MSO_LANGUAGE_ID.RUSSIAN
    for paragraph in tf.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER


# функция для выравнивания и изменения размеров надписи (placeholders)
def size_change_plh(words, left, top, width, height):
    words.left = Cm(left)
    words.top = Cm(top)
    words.width = Cm(width)
    words.height = Cm(height)


# добавить пробел
def add_space(placeholders):
    tf = placeholders.text_frame
    tf.add_paragraph()  # пробел между текстом на слайде


# функция для добавления фигуры
def pict_for_background(l, t, w, h, name_picture):
    pic = slide.shapes.add_picture(name_picture, left=Cm(l), top=Cm(t), height=Cm(h), width=Cm(w))

    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


root = Presentation()

# титульный слайд
first_slide_layout = root.slide_layouts[0]
slide = root.slides.add_slide(first_slide_layout)

title_placeholder = slide.shapes.title
title_placeholder.text = 'Введение в Data Science и машинное обучение'
size_change_plh(title_placeholder, 3.81, 3.92, 17.78, 3.49)

presentation_title(title_placeholder)

subtitle = slide.placeholders[1]
subtitle.text = 'Автор: Талейко Ирина'
size_change_plh(subtitle, 6.2, 8.14, 13, 1.42)

text_frame_theory = subtitle.text_frame
for run in text_frame_theory.paragraphs[0].runs:
    run.font.size = Pt(27)
    run.font.bold = False
    run.font.name = "Book Antiqua"

pic = slide.shapes.add_picture('public-speaker-givin.png', Inches(4.157), Inches(3.752), Inches(5.748), Inches(3.484))

rectangle = slide.shapes.add_shape(autoshape_type_id=MSO_SHAPE.RECTANGLE,
                                   left=Cm(0),
                                   top=Cm(0.54),
                                   height=Cm(17.97),
                                   width=Cm(25.4))
rectangleFill = rectangle.fill
rectangleFill.solid()
rectangleFillColour = rectangleFill.fore_color
rectangleFillColour.rgb = RGBColor(255,255,255)

slide.shapes._spTree.remove(rectangle._element)
slide.shapes._spTree.insert(2, rectangle._element)

# следующий слайд
slide = root.slides.add_slide(root.slide_layouts[0])
shapes = slide.shapes.title
shapes.text = 'Основы статистики'
presentation_title(shapes)
pict_for_background(0, 0, 25.4, 19.05, 'fon.png')
theory = slide.placeholders[1]
add_space(theory)

# следующий слайд
slide = root.slides.add_slide(root.slide_layouts[1])
pict_for_background(0, 0, 25.4, 19.05, 'fon.png')
shapes = slide.shapes.title.text = 'Основы статистики'
title_format(shapes)

theory = slide.placeholders[1]

paragraph = "Базовые понятия:"
format_placeholders(theory, "Базовые понятия:")

add_indent_paragraph(theory, "выборка")
add_indent_paragraph(theory, "генеральная совокупность")
add_indent_paragraph(theory, "статистически значимые различия")

add_paragraph_format(theory, "Методы:")

add_indent_paragraph(theory, "регрессионный анализ")
add_indent_paragraph(theory, "дисперсионный анализ")
add_indent_paragraph(theory, "Т-критерий")
add_indent_paragraph(theory, "коэффициенты корреляции")

# следующий слайд
slide = root.slides.add_slide(root.slide_layouts[1])
pict_for_background(0, 0, 25.4, 19.05, 'fon.png')
shapes = slide.shapes.title.text = 'Основы статистики'
title_format(shapes)

theory = slide.placeholders[1]

format_placeholders(theory, "Способы выборки:")

add_indent_paragraph(theory, "простая случайная выборка")
add_indent_paragraph(theory, "стратифицированная выборка (деление генеральной совокупности на страты)")
add_indent_paragraph(theory, "групповая выборка")

add_paragraph_format(theory, "Типы переменных:")

add_indent_paragraph(theory, "количественные")
add_indent_paragraph(theory, "качественные")
add_indent_paragraph(theory, "ранговые")

# следующий слайд
slide = root.slides.add_slide(root.slide_layouts[1])
pict_for_background(0, 0, 25.4, 19.05, 'fon.png')
shapes = slide.shapes.title.text = 'Меры центральной тенденции'
title_format(shapes)

theory = slide.placeholders[1]

format_placeholders(theory, "Мода - значение измеряемого признака, которое встречается максимально часто (мода распределения).")
add_space(theory)

add_paragraph_format(theory, "Медиана - значение признака, которое делит упорядоченное множество пополам.")
add_space(theory)

add_paragraph_format(theory, "Среднее значение - сумма всех значений, деленная на количество измеренных значений.")

# следующий слайд
slide = root.slides.add_slide(root.slide_layouts[5])
pict_for_background(0, 0, 25.4, 19.05, 'fon.png')
slide.shapes.title.text = 'Нормальное распределение'
shapes = slide.shapes.title.text_frame
title_format(shapes)

paragraph_centered("Является унимодальным и симметричным", 3.57, 4.01, 18.27, 1.2)
paragraph_centered("Отклонения наблюдений от среднего подчиняются определенному вероятностному закону", 3.16, 16.17, 19.08, 2.14)

# График нормального распределения
mu = 0      # Среднее
sigma = 1   # Стандартное отклонение
# Массив значений по оси X
x = np.linspace(mu - 4*sigma, mu + 4*sigma, 1000)
# Значения плотности вероятности
pdf = norm.pdf(x, mu, sigma)

# Параметры графика
plt.figure(figsize=(10, 6))
plt.plot(x, pdf, label=f'N({mu}, {sigma**2})', color='blue')
plt.title('График нормального распределения')
plt.xlabel('X-axis')
plt.ylabel('Плотность вероятности')
plt.legend()
plt.grid()

plt.savefig('normal_distribution.png')
plt.close()

img_path = 'normal_distribution.png'
slide.shapes.add_picture(img_path, Inches(2.028), Inches(2.083), Inches(5.949), Inches(4.161))

# следующий слайд
slide = root.slides.add_slide(root.slide_layouts[5])
pict_for_background(0, 0, 25.4, 19.05, 'fon.png')
shapes = slide.shapes.title.text = 'Меры изменчивости'
title_format(shapes)

add_placeholder("Размах - разница между минимальным и максимальным значением выборки.", 1.27, 4.45, 22.86, 6)

theory = slide.shapes[2]
add_space(theory)

add_paragraph_format(theory, "Дисперсия - средний квадрат отклонений индивидуальных значений признака от их средней величины.")

slide.shapes.add_picture('дисперсия формула.png', Inches(3.295), Inches(4.378), Inches(3.26), Inches(1.567))

add_placeholder("n - число значений", 17.03, 13.9, 7.48, 1.2)

add_placeholder("Квадратный корень дисперсии – среднее квадратное отклонение, Ꝺ = D**0,5.", 1.27, 16.19, 22.86, 2.31)

# следующий слайд
slide = root.slides.add_slide(root.slide_layouts[6])
pict_for_background(0, 0, 25.4, 19.05, 'fon.png')
paragraph = "Если среднее квадратическое отклонение было найдено не для всей генеральной совокупности, а только для конкретной выборки, то показатель называется SD - стандартное отклонение."
paragraph_centered(paragraph, 1.27, 4.21, 22.86, 6)

slide.shapes.add_picture('дисперсия выборки формула.png', Inches(3.295), Inches(4.378), Inches(3.26), Inches(1.567))

# следующий слайд
slide = root.slides.add_slide(root.slide_layouts[1])
pict_for_background(0, 0, 25.4, 19.05, 'fon.png')
shapes = slide.shapes.title.text = 'Свойства дисперсии'
title_format(shapes)

theory = slide.placeholders[1]

paragraph = "Если к каждому значению выборки прибавить определенное число, то ни дисперсия, ни стандартное отклонение не изменятся."
format_placeholders(theory, paragraph)
add_space(theory)
thesis = "Если мы каждое значение выборки умножим на определенное число, то стандартное отклонение увеличится в это число раз, а дисперсия увеличится в это число раз в квадрате."
add_paragraph_format(theory, thesis)

slide.shapes.add_picture('изм. sd.png', Inches(0.9016), Inches(5.756), Inches(3.752), Inches(1.024))
slide.shapes.add_picture('изм. дисперсии.png', Inches(5.118), Inches(5.756), Inches(3.752), Inches(1.024))

# заливка фона слайдов
for slide in root.slides:
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(145, 100, 100)


root.save('Basics of Statistics.pptx')

